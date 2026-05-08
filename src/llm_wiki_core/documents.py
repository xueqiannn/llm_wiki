from __future__ import annotations

from pathlib import Path


TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".csv", ".json", ".log"}


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)

    return path.read_bytes().decode("utf-8", errors="replace")


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(page.strip() for page in pages if page.strip())


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text.strip():
                slide_text.append(text.strip())
        if slide_text:
            chunks.append(f"## Slide {slide_no}\n" + "\n".join(slide_text))
    return "\n\n".join(chunks)
